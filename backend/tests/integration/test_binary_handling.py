"""
Quick test to verify binary file fix
"""
import asyncio
import sys
import base64

sys.path.insert(0, '.')

from code_executor import code_executor, ExecutionStatus

async def test_excel():
    print("Testing Excel file generation...")

    code = """
import pandas as pd
df = pd.DataFrame({
    'Nome': ['Alice', 'Bob', 'Carol'],
    'Nota': [95, 87, 92]
})
df.to_excel('test.xlsx', index=False)
print('Excel file created!')
"""

    result = await code_executor.execute(
        code=code,
        libraries=["pandas", "openpyxl"],
        output_files=["test.xlsx"]
    )

    print(f"Status: {result.status}")
    print(f"stdout: {result.stdout}")
    print(f"stderr: {result.stderr}")
    print(f"Files generated: {len(result.files_generated)}")

    if result.files_generated:
        file = result.files_generated[0]
        print(f"  - {file.filename}: {file.size_bytes} bytes, {file.mime_type}")

        # Decode and check the file signature (Excel files start with PK)
        data = base64.b64decode(file.content_base64)
        print(f"  - First 4 bytes: {data[:4]}")

        # XLSX files are ZIP files, they start with PK (0x50 0x4B)
        if data[:2] == b'PK':
            print("  - [OK] Valid XLSX file signature (PK/ZIP)")
        else:
            print(f"  - [FAIL] Invalid signature, expected PK, got {data[:2]}")

        # Save to verify
        with open('test_output.xlsx', 'wb') as f:
            f.write(data)
        print("  - Saved to test_output.xlsx - try opening it!")
    else:
        print("[FAIL] No files generated")

async def test_pdf():
    print("\nTesting PDF file generation...")

    code = """
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

c = canvas.Canvas('test.pdf', pagesize=letter)
c.drawString(100, 750, 'Hello, PDF World!')
c.save()
print('PDF created!')
"""

    result = await code_executor.execute(
        code=code,
        libraries=["reportlab"],
        output_files=["test.pdf"]
    )

    print(f"Status: {result.status}")
    print(f"stdout: {result.stdout}")
    print(f"stderr: {result.stderr}")
    print(f"error_message: {result.error_message}")
    print(f"Files generated: {len(result.files_generated)}")

    if result.files_generated:
        file = result.files_generated[0]
        print(f"  - {file.filename}: {file.size_bytes} bytes")

        data = base64.b64decode(file.content_base64)
        print(f"  - First 8 bytes: {data[:8]}")

        # PDF files start with %PDF
        if data[:4] == b'%PDF':
            print("  - [OK] Valid PDF file signature")
        else:
            print(f"  - [FAIL] Invalid signature, expected %PDF, got {data[:4]}")

        with open('test_output.pdf', 'wb') as f:
            f.write(data)
        print("  - Saved to test_output.pdf")

async def test_pptx():
    print("\nTesting PowerPoint file generation...")

    code = """
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

left = Inches(2)
top = Inches(2)
width = Inches(4)
height = Inches(1)

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Hello, PowerPoint!"

prs.save('test.pptx')
print('PPTX created!')
"""

    result = await code_executor.execute(
        code=code,
        libraries=["python-pptx"],
        output_files=["test.pptx"]
    )

    print(f"Status: {result.status}")
    print(f"stdout: {result.stdout}")
    print(f"stderr: {result.stderr}")
    print(f"error_message: {result.error_message}")
    print(f"Files generated: {len(result.files_generated)}")

    if result.files_generated:
        file = result.files_generated[0]
        print(f"  - {file.filename}: {file.size_bytes} bytes")

        data = base64.b64decode(file.content_base64)
        print(f"  - First 4 bytes: {data[:4]}")

        # PPTX files are ZIP files, they start with PK
        if data[:2] == b'PK':
            print("  - [OK] Valid PPTX file signature (PK/ZIP)")
        else:
            print(f"  - [FAIL] Invalid signature, expected PK, got {data[:2]}")

        with open('test_output.pptx', 'wb') as f:
            f.write(data)
        print("  - Saved to test_output.pptx")

async def main():
    print("="*50)
    print(" BINARY FILE FIX VERIFICATION")
    print("="*50)

    await test_excel()
    await test_pdf()
    await test_pptx()

    print("\n" + "="*50)
    print(" DONE - Check the test_output.* files!")
    print("="*50)

if __name__ == "__main__":
    asyncio.run(main())
