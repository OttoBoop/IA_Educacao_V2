"""
PROVA AI - Sistema de Anexos v2.2

Envia arquivos NATIVAMENTE para APIs multimodais.
NÃO converte nada - envia o arquivo original!

Formatos suportados por provider:
- OpenAI (GPT-4o): PDF, imagens (PNG, JPG, GIF, WEBP)
- Anthropic (Claude): PDF, imagens (PNG, JPG, GIF, WEBP)
- Google (Gemini): PDF, imagens, vídeo, áudio

Para arquivos de texto/código, envia como texto no prompt.
"""

import base64
import mimetypes
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime
import json
import httpx


# ============================================================
# TIPOS DE ARQUIVO SUPORTADOS
# ============================================================

# Arquivos que são enviados como ANEXO BINÁRIO (base64)
FORMATOS_BINARIOS = {
    # PDFs
    '.pdf': 'application/pdf',
    
    # Imagens
    '.png': 'image/png',
    '.jpg': 'image/jpeg',
    '.jpeg': 'image/jpeg',
    '.gif': 'image/gif',
    '.webp': 'image/webp',
    '.bmp': 'image/bmp',
    '.tiff': 'image/tiff',
    '.tif': 'image/tiff',
}

# Arquivos que são enviados como TEXTO no prompt
FORMATOS_TEXTO = {
    # Documentos
    '.txt': 'text/plain',
    '.md': 'text/markdown',
    '.csv': 'text/csv',
    '.json': 'application/json',
    '.xml': 'application/xml',
    '.html': 'text/html',
    '.htm': 'text/html',
    
    # Código
    '.py': 'text/x-python',
    '.js': 'text/javascript',
    '.ts': 'text/typescript',
    '.java': 'text/x-java',
    '.c': 'text/x-c',
    '.cpp': 'text/x-c++',
    '.h': 'text/x-c',
    '.cs': 'text/x-csharp',
    '.rb': 'text/x-ruby',
    '.go': 'text/x-go',
    '.rs': 'text/x-rust',
    '.php': 'text/x-php',
    '.sql': 'text/x-sql',
    '.r': 'text/x-r',
    '.m': 'text/x-matlab',
    '.sh': 'text/x-shellscript',
    '.bash': 'text/x-shellscript',
    '.ps1': 'text/x-powershell',
    '.yaml': 'text/yaml',
    '.yml': 'text/yaml',
    '.toml': 'text/toml',
    '.ini': 'text/ini',
    '.cfg': 'text/plain',
    '.conf': 'text/plain',
    '.log': 'text/plain',
    
    # Notebooks
    '.ipynb': 'application/json',  # Jupyter notebooks são JSON
    
    # LaTeX
    '.tex': 'text/x-latex',
    '.latex': 'text/x-latex',
}

# Arquivos que precisam de conversão especial
FORMATOS_ESPECIAIS = {
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.xls': 'application/vnd.ms-excel',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
}

# Modelos que usam reasoning e não suportam certos parâmetros
# Deprecated: 'o1', 'o1-pro' (removidos pois estão deprecated)
REASONING_MODELS = ['o3', 'o3-mini', 'o3-pro', 'o4-mini', 'gpt-5', 'gpt-5-mini', 'gpt-5-nano', 'gpt-5.1', 'gpt-5.2', 'deepseek-reasoner']


def is_reasoning_model(modelo: str) -> bool:
    """Verifica se o modelo é de reasoning (não suporta temperature)"""
    if not modelo:
        return False
    modelo_lower = modelo.lower()
    return any(r in modelo_lower for r in REASONING_MODELS)


@dataclass
class ArquivoAnexo:
    """Representa um arquivo preparado para envio"""
    nome: str
    caminho: str
    extensao: str
    mime_type: str
    tamanho_bytes: int
    
    # Conteúdo
    conteudo_base64: Optional[str] = None  # Para binários
    conteudo_texto: Optional[str] = None   # Para texto
    
    # Metadados
    tipo_envio: str = "binario"  # "binario", "texto", "especial"
    suportado: bool = True
    aviso: Optional[str] = None
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "nome": self.nome,
            "extensao": self.extensao,
            "mime_type": self.mime_type,
            "tamanho_bytes": self.tamanho_bytes,
            "tipo_envio": self.tipo_envio,
            "suportado": self.suportado,
            "aviso": self.aviso,
            "tem_conteudo": bool(self.conteudo_base64 or self.conteudo_texto)
        }


class PreparadorArquivos:
    """Prepara arquivos para envio às APIs"""
    
    def __init__(self, max_tamanho_mb: float = 20.0):
        self.max_tamanho_bytes = int(max_tamanho_mb * 1024 * 1024)
    
    def preparar(self, caminho: str) -> ArquivoAnexo:
        """
        Prepara um arquivo para envio.
        
        - Binários (PDF, imagens): converte para base64
        - Texto (código, CSV, etc.): lê como texto
        - Especiais (DOCX, XLSX): tenta extrair texto
        """
        arquivo = Path(caminho)
        
        if not arquivo.exists():
            raise FileNotFoundError(f"Arquivo não encontrado: {caminho}")
        
        extensao = arquivo.suffix.lower()
        nome = arquivo.name
        tamanho = arquivo.stat().st_size
        
        # Verificar tamanho
        if tamanho > self.max_tamanho_bytes:
            return ArquivoAnexo(
                nome=nome,
                caminho=str(arquivo),
                extensao=extensao,
                mime_type="application/octet-stream",
                tamanho_bytes=tamanho,
                tipo_envio="erro",
                suportado=False,
                aviso=f"Arquivo muito grande ({tamanho / 1024 / 1024:.1f} MB). Máximo: {self.max_tamanho_bytes / 1024 / 1024:.1f} MB"
            )
        
        # Determinar tipo de envio
        if extensao in FORMATOS_BINARIOS:
            return self._preparar_binario(arquivo, extensao)
        elif extensao in FORMATOS_TEXTO:
            return self._preparar_texto(arquivo, extensao)
        elif extensao in FORMATOS_ESPECIAIS:
            return self._preparar_especial(arquivo, extensao)
        else:
            # Tentar como texto
            return self._tentar_como_texto(arquivo, extensao)
    
    def _preparar_binario(self, arquivo: Path, extensao: str) -> ArquivoAnexo:
        """Prepara arquivo binário (PDF, imagem) como base64"""
        mime_type = FORMATOS_BINARIOS[extensao]
        
        with open(arquivo, 'rb') as f:
            conteudo = f.read()
        
        conteudo_base64 = base64.standard_b64encode(conteudo).decode('utf-8')
        
        return ArquivoAnexo(
            nome=arquivo.name,
            caminho=str(arquivo),
            extensao=extensao,
            mime_type=mime_type,
            tamanho_bytes=len(conteudo),
            conteudo_base64=conteudo_base64,
            tipo_envio="binario",
            suportado=True
        )
    
    def _preparar_texto(self, arquivo: Path, extensao: str) -> ArquivoAnexo:
        """Prepara arquivo de texto"""
        mime_type = FORMATOS_TEXTO.get(extensao, 'text/plain')
        
        # Tentar ler com diferentes encodings
        conteudo = None
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                with open(arquivo, 'r', encoding=encoding) as f:
                    conteudo = f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if conteudo is None:
            return ArquivoAnexo(
                nome=arquivo.name,
                caminho=str(arquivo),
                extensao=extensao,
                mime_type=mime_type,
                tamanho_bytes=arquivo.stat().st_size,
                tipo_envio="erro",
                suportado=False,
                aviso="Não foi possível ler o arquivo (encoding desconhecido)"
            )
        
        return ArquivoAnexo(
            nome=arquivo.name,
            caminho=str(arquivo),
            extensao=extensao,
            mime_type=mime_type,
            tamanho_bytes=len(conteudo.encode('utf-8')),
            conteudo_texto=conteudo,
            tipo_envio="texto",
            suportado=True
        )
    
    def _preparar_especial(self, arquivo: Path, extensao: str) -> ArquivoAnexo:
        """Prepara arquivos especiais (DOCX, XLSX)"""
        mime_type = FORMATOS_ESPECIAIS[extensao]
        tamanho = arquivo.stat().st_size
        
        # Tentar extrair texto
        conteudo_texto = None
        aviso = None
        
        if extensao == '.docx':
            conteudo_texto, aviso = self._extrair_texto_docx(arquivo)
        elif extensao in ['.xlsx', '.xls']:
            conteudo_texto, aviso = self._extrair_texto_excel(arquivo)
        elif extensao == '.pptx':
            conteudo_texto, aviso = self._extrair_texto_pptx(arquivo)
        
        if conteudo_texto:
            return ArquivoAnexo(
                nome=arquivo.name,
                caminho=str(arquivo),
                extensao=extensao,
                mime_type=mime_type,
                tamanho_bytes=tamanho,
                conteudo_texto=conteudo_texto,
                tipo_envio="texto_extraido",
                suportado=True,
                aviso=aviso
            )
        else:
            # Fallback: enviar como binário (algumas APIs aceitam)
            with open(arquivo, 'rb') as f:
                conteudo = f.read()
            
            return ArquivoAnexo(
                nome=arquivo.name,
                caminho=str(arquivo),
                extensao=extensao,
                mime_type=mime_type,
                tamanho_bytes=tamanho,
                conteudo_base64=base64.standard_b64encode(conteudo).decode('utf-8'),
                tipo_envio="binario",
                suportado=True,
                aviso=aviso or "Enviado como binário (extração de texto falhou)"
            )
    
    def _extrair_texto_docx(self, arquivo: Path) -> Tuple[Optional[str], Optional[str]]:
        """Extrai texto de DOCX"""
        try:
            from docx import Document
            doc = Document(arquivo)
            paragrafos = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragrafos), None
        except ImportError:
            return None, "python-docx não instalado. Instale com: pip install python-docx"
        except Exception as e:
            return None, f"Erro ao extrair texto: {str(e)}"
    
    def _extrair_texto_excel(self, arquivo: Path) -> Tuple[Optional[str], Optional[str]]:
        """Extrai texto de Excel"""
        try:
            import pandas as pd
            
            # Ler todas as sheets
            xl = pd.ExcelFile(arquivo)
            textos = []
            
            for sheet_name in xl.sheet_names:
                df = pd.read_excel(xl, sheet_name=sheet_name)
                textos.append(f"=== Sheet: {sheet_name} ===")
                textos.append(df.to_string())
            
            return "\n\n".join(textos), None
        except ImportError:
            return None, "pandas/openpyxl não instalado. Instale com: pip install pandas openpyxl"
        except Exception as e:
            return None, f"Erro ao extrair texto: {str(e)}"
    
    def _extrair_texto_pptx(self, arquivo: Path) -> Tuple[Optional[str], Optional[str]]:
        """Extrai texto de PowerPoint"""
        try:
            from pptx import Presentation
            prs = Presentation(arquivo)
            
            textos = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"=== Slide {i} ==="]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                textos.append("\n".join(slide_text))
            
            return "\n\n".join(textos), None
        except ImportError:
            return None, "python-pptx não instalado. Instale com: pip install python-pptx"
        except Exception as e:
            return None, f"Erro ao extrair texto: {str(e)}"
    
    def _tentar_como_texto(self, arquivo: Path, extensao: str) -> ArquivoAnexo:
        """Tenta ler arquivo desconhecido como texto"""
        try:
            with open(arquivo, 'r', encoding='utf-8') as f:
                conteudo = f.read()
            
            return ArquivoAnexo(
                nome=arquivo.name,
                caminho=str(arquivo),
                extensao=extensao,
                mime_type='text/plain',
                tamanho_bytes=len(conteudo.encode('utf-8')),
                conteudo_texto=conteudo,
                tipo_envio="texto",
                suportado=True,
                aviso=f"Formato {extensao} não reconhecido, tratado como texto"
            )
        except:
            # Último recurso: binário
            with open(arquivo, 'rb') as f:
                conteudo = f.read()
            
            return ArquivoAnexo(
                nome=arquivo.name,
                caminho=str(arquivo),
                extensao=extensao,
                mime_type='application/octet-stream',
                tamanho_bytes=len(conteudo),
                conteudo_base64=base64.standard_b64encode(conteudo).decode('utf-8'),
                tipo_envio="binario",
                suportado=True,
                aviso=f"Formato {extensao} não reconhecido, enviado como binário"
            )


# ============================================================
# CLIENTE DE API COM VERIFICAÇÃO DE ENVIO
# ============================================================

@dataclass
class ResultadoEnvio:
    """Resultado do envio de mensagem com anexos"""
    sucesso: bool
    resposta: str = ""
    
    # Metadados
    provider: str = ""
    modelo: str = ""
    tokens_entrada: int = 0
    tokens_saida: int = 0
    
    # Verificação de anexos
    anexos_enviados: List[Dict[str, Any]] = field(default_factory=list)
    anexos_confirmados: bool = False
    
    # Erros
    erro: Optional[str] = None
    erro_detalhes: Optional[str] = None
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "sucesso": self.sucesso,
            "resposta": self.resposta,
            "provider": self.provider,
            "modelo": self.modelo,
            "tokens_entrada": self.tokens_entrada,
            "tokens_saida": self.tokens_saida,
            "anexos_enviados": self.anexos_enviados,
            "anexos_confirmados": self.anexos_confirmados,
            "erro": self.erro,
            "erro_detalhes": self.erro_detalhes
        }


class ClienteAPIMultimodal:
    """
    Cliente para APIs multimodais com verificação de envio.
    
    Suporta:
    - OpenAI (GPT-4o, GPT-4-turbo)
    - Anthropic (Claude 3.x)
    - Google (Gemini 1.5)
    """
    
    def __init__(self, config: Dict[str, Any]):
        self.tipo = config.get("tipo", "openai")
        self.api_key = config.get("api_key", "")
        self.modelo = config.get("modelo", "gpt-4o")
        self.base_url = config.get("base_url")
        self.max_tokens = config.get("max_tokens", 4096)
        self.temperature = config.get("temperature", 0.7)
        # Check if model supports temperature (reasoning models don't)
        self.suporta_temperature = config.get("suporta_temperature", not is_reasoning_model(self.modelo))

        self.preparador = PreparadorArquivos()
    
    async def enviar_com_anexos(
        self,
        mensagem: str,
        arquivos: List[str],  # Lista de caminhos de arquivo
        system_prompt: str = None,
        historico: List[Dict] = None,
        verificar_anexos: bool = True
    ) -> ResultadoEnvio:
        """
        Envia mensagem com arquivos anexados.
        
        Args:
            mensagem: Texto da mensagem
            arquivos: Lista de caminhos de arquivos para anexar
            system_prompt: Prompt de sistema
            historico: Mensagens anteriores
            verificar_anexos: Se deve verificar se os anexos foram recebidos
        
        Returns:
            ResultadoEnvio com resposta e status dos anexos
        """
        # Preparar arquivos
        anexos_preparados = []
        for caminho in arquivos:
            try:
                anexo = self.preparador.preparar(caminho)
                anexos_preparados.append(anexo)
            except Exception as e:
                return ResultadoEnvio(
                    sucesso=False,
                    erro=f"Erro ao preparar arquivo {caminho}",
                    erro_detalhes=str(e)
                )
        
        # Verificar se há anexos não suportados
        nao_suportados = [a for a in anexos_preparados if not a.suportado]
        if nao_suportados:
            return ResultadoEnvio(
                sucesso=False,
                erro="Arquivos não suportados",
                erro_detalhes="; ".join([f"{a.nome}: {a.aviso}" for a in nao_suportados]),
                anexos_enviados=[a.to_dict() for a in anexos_preparados]
            )
        
        # Enviar para API apropriada
        try:
            if self.tipo in ["openai", "openrouter"]:
                resultado = await self._enviar_openai(mensagem, anexos_preparados, system_prompt, historico)
            elif self.tipo == "anthropic":
                resultado = await self._enviar_anthropic(mensagem, anexos_preparados, system_prompt, historico)
            elif self.tipo == "google":
                resultado = await self._enviar_google(mensagem, anexos_preparados, system_prompt, historico)
            else:
                # Fallback: enviar textos, ignorar binários
                resultado = await self._enviar_texto_apenas(mensagem, anexos_preparados, system_prompt, historico)
            
            # Adicionar info dos anexos
            resultado.anexos_enviados = [a.to_dict() for a in anexos_preparados]
            
            # Verificar se a IA confirmou receber os anexos
            if verificar_anexos and resultado.sucesso:
                resultado.anexos_confirmados = self._verificar_confirmacao_anexos(
                    resultado.resposta,
                    anexos_preparados
                )
            
            return resultado
            
        except Exception as e:
            return ResultadoEnvio(
                sucesso=False,
                erro="Erro ao enviar para API",
                erro_detalhes=str(e),
                anexos_enviados=[a.to_dict() for a in anexos_preparados]
            )
    
    async def _enviar_openai(
        self,
        mensagem: str,
        anexos: List[ArquivoAnexo],
        system_prompt: str,
        historico: List[Dict]
    ) -> ResultadoEnvio:
        """Envia para OpenAI com anexos"""

        # Check if this is a reasoning model (o1, o3, etc.)
        is_reasoning = is_reasoning_model(self.modelo)

        # Reasoning models don't support vision/image content
        if is_reasoning:
            has_images = any(
                a.tipo_envio == "binario" and a.mime_type.startswith('image/')
                for a in anexos
            )
            if has_images:
                return ResultadoEnvio(
                    sucesso=False,
                    erro=f"Modelo {self.modelo} não suporta análise de imagens",
                    erro_detalhes="Modelos de reasoning (o1, o3, etc.) não suportam conteúdo visual. Use um modelo com suporte a vision como gpt-4o.",
                    provider="openai",
                    modelo=self.modelo
                )

        messages = []

        # Reasoning models use "developer" role instead of "system"
        if system_prompt:
            role = "developer" if is_reasoning else "system"
            messages.append({"role": role, "content": system_prompt})

        if historico:
            for msg in historico:
                role = msg.get("role", "user")
                # Convert system to developer for reasoning models
                if role == "system" and is_reasoning:
                    role = "developer"
                messages.append({"role": role, "content": msg.get("content", "")})

        # Construir mensagem do usuário com anexos
        content = []

        for anexo in anexos:
            if anexo.tipo_envio == "binario" and anexo.conteudo_base64:
                if anexo.extensao == '.pdf':
                    # OpenAI suporta PDF nativamente (desde GPT-4o)
                    content.append({
                        "type": "file",
                        "file": {
                            "filename": anexo.nome,
                            "file_data": f"data:{anexo.mime_type};base64,{anexo.conteudo_base64}"
                        }
                    })
                elif anexo.mime_type.startswith('image/'):
                    content.append({
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{anexo.mime_type};base64,{anexo.conteudo_base64}",
                            "detail": "high"
                        }
                    })

            elif anexo.conteudo_texto:
                # Arquivos de texto vão como texto
                content.append({
                    "type": "text",
                    "text": f"--- ARQUIVO: {anexo.nome} ---\n{anexo.conteudo_texto}\n--- FIM ARQUIVO ---"
                })

        # Adicionar mensagem do usuário
        content.append({"type": "text", "text": mensagem})

        messages.append({"role": "user", "content": content})

        # Fazer requisição
        url = self.base_url or "https://api.openai.com/v1"
        if not url.endswith("/chat/completions"):
            url = f"{url.rstrip('/')}/chat/completions"

        # Build params based on model type
        params = {
            "model": self.modelo,
            "messages": messages,
        }

        # Reasoning models use max_completion_tokens, others use max_tokens
        if is_reasoning:
            params["max_completion_tokens"] = self.max_tokens
        else:
            params["max_tokens"] = self.max_tokens

        # Only add temperature if model supports it
        if not is_reasoning and self.suporta_temperature and self.temperature is not None:
            params["temperature"] = self.temperature

        async with httpx.AsyncClient(timeout=180.0) as client:
            response = await client.post(
                url,
                headers={
                    "Authorization": f"Bearer {self.api_key}",
                    "Content-Type": "application/json"
                },
                json=params
            )

            if response.status_code != 200:
                return ResultadoEnvio(
                    sucesso=False,
                    erro=f"Erro API OpenAI: {response.status_code}",
                    erro_detalhes=response.text,
                    provider="openai",
                    modelo=self.modelo
                )

            data = response.json()

            return ResultadoEnvio(
                sucesso=True,
                resposta=data["choices"][0]["message"]["content"],
                provider="openai",
                modelo=self.modelo,
                tokens_entrada=data.get("usage", {}).get("prompt_tokens", 0),
                tokens_saida=data.get("usage", {}).get("completion_tokens", 0)
            )
    
    async def _enviar_anthropic(
        self,
        mensagem: str,
        anexos: List[ArquivoAnexo],
        system_prompt: str,
        historico: List[Dict]
    ) -> ResultadoEnvio:
        """Envia para Anthropic com anexos"""
        
        messages = []
        
        if historico:
            messages.extend(historico)
        
        # Construir mensagem com anexos
        content = []
        
        for anexo in anexos:
            if anexo.tipo_envio == "binario" and anexo.conteudo_base64:
                if anexo.extensao == '.pdf':
                    # Claude suporta PDF nativamente
                    content.append({
                        "type": "document",
                        "source": {
                            "type": "base64",
                            "media_type": anexo.mime_type,
                            "data": anexo.conteudo_base64
                        }
                    })
                elif anexo.mime_type.startswith('image/'):
                    content.append({
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": anexo.mime_type,
                            "data": anexo.conteudo_base64
                        }
                    })
            
            elif anexo.conteudo_texto:
                content.append({
                    "type": "text",
                    "text": f"--- ARQUIVO: {anexo.nome} ---\n{anexo.conteudo_texto}\n--- FIM ARQUIVO ---"
                })
        
        content.append({"type": "text", "text": mensagem})
        
        messages.append({"role": "user", "content": content})
        
        # Fazer requisição
        url = self.base_url or "https://api.anthropic.com/v1/messages"

        # Build params
        params = {
            "model": self.modelo,
            "max_tokens": self.max_tokens,
            "system": system_prompt or "Você é um assistente útil.",
            "messages": messages
        }

        # Add temperature if model supports it
        if self.suporta_temperature and self.temperature is not None:
            params["temperature"] = self.temperature

        async with httpx.AsyncClient(timeout=180.0) as client:
            response = await client.post(
                url,
                headers={
                    "x-api-key": self.api_key,
                    "anthropic-version": "2023-06-01",
                    "Content-Type": "application/json"
                },
                json=params
            )
            
            if response.status_code != 200:
                return ResultadoEnvio(
                    sucesso=False,
                    erro=f"Erro API Anthropic: {response.status_code}",
                    erro_detalhes=response.text,
                    provider="anthropic",
                    modelo=self.modelo
                )
            
            data = response.json()
            
            resposta_texto = ""
            for block in data.get("content", []):
                if block.get("type") == "text":
                    resposta_texto += block.get("text", "")
            
            return ResultadoEnvio(
                sucesso=True,
                resposta=resposta_texto,
                provider="anthropic",
                modelo=self.modelo,
                tokens_entrada=data.get("usage", {}).get("input_tokens", 0),
                tokens_saida=data.get("usage", {}).get("output_tokens", 0)
            )
    
    async def _enviar_google(
        self,
        mensagem: str,
        anexos: List[ArquivoAnexo],
        system_prompt: str,
        historico: List[Dict]
    ) -> ResultadoEnvio:
        """Envia para Google Gemini com anexos"""

        parts = []

        # Add file attachments
        for anexo in anexos:
            if anexo.tipo_envio == "binario" and anexo.conteudo_base64:
                parts.append({
                    "inline_data": {
                        "mime_type": anexo.mime_type,
                        "data": anexo.conteudo_base64
                    }
                })
                parts.append({"text": f"[Arquivo acima: {anexo.nome}]"})
            
            elif anexo.conteudo_texto:
                parts.append({"text": f"--- ARQUIVO: {anexo.nome} ---\n{anexo.conteudo_texto}\n--- FIM ---"})
        
        parts.append({"text": mensagem})
        
        contents = [{"role": "user", "parts": parts}]
        
        # Fazer requisição
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{self.modelo}:generateContent"

        # Build generation config
        generation_config = {
            "maxOutputTokens": self.max_tokens,
        }

        # Add temperature if model supports it
        if self.suporta_temperature and self.temperature is not None:
            generation_config["temperature"] = self.temperature

        # Build request body
        request_body = {
            "contents": contents,
            "generationConfig": generation_config
        }

        # Add system instruction if provided
        if system_prompt:
            request_body["system_instruction"] = {
                "parts": [{"text": system_prompt}]
            }

        async with httpx.AsyncClient(timeout=180.0) as client:
            response = await client.post(
                url,
                params={"key": self.api_key},
                headers={"Content-Type": "application/json"},
                json=request_body
            )
            
            if response.status_code != 200:
                return ResultadoEnvio(
                    sucesso=False,
                    erro=f"Erro API Google: {response.status_code}",
                    erro_detalhes=response.text,
                    provider="google",
                    modelo=self.modelo
                )
            
            data = response.json()
            
            resposta_texto = ""
            for candidate in data.get("candidates", []):
                for part in candidate.get("content", {}).get("parts", []):
                    resposta_texto += part.get("text", "")
            
            return ResultadoEnvio(
                sucesso=True,
                resposta=resposta_texto,
                provider="google",
                modelo=self.modelo,
                tokens_entrada=data.get("usageMetadata", {}).get("promptTokenCount", 0),
                tokens_saida=data.get("usageMetadata", {}).get("candidatesTokenCount", 0)
            )
    
    async def _enviar_texto_apenas(
        self,
        mensagem: str,
        anexos: List[ArquivoAnexo],
        system_prompt: str,
        historico: List[Dict]
    ) -> ResultadoEnvio:
        """Fallback: envia apenas textos para providers sem suporte a binários"""
        
        # Construir prompt com textos
        textos = []
        binarios_ignorados = []
        
        for anexo in anexos:
            if anexo.conteudo_texto:
                textos.append(f"--- ARQUIVO: {anexo.nome} ---\n{anexo.conteudo_texto}\n---")
            else:
                binarios_ignorados.append(anexo.nome)
        
        prompt_completo = mensagem
        if textos:
            prompt_completo = "\n\n".join(textos) + "\n\n" + mensagem
        
        aviso = None
        if binarios_ignorados:
            aviso = f"Arquivos binários ignorados (provider não suporta): {', '.join(binarios_ignorados)}"
        
        # Usar OpenAI como fallback básico
        return await self._enviar_openai(prompt_completo, [], system_prompt, historico)
    
    def _verificar_confirmacao_anexos(self, resposta: str, anexos: List[ArquivoAnexo]) -> bool:
        """
        Verifica se a IA aparenta ter recebido/processado os anexos.
        
        Procura por indicações de que a IA viu o conteúdo dos arquivos.
        """
        resposta_lower = resposta.lower()
        
        # Indicadores positivos
        indicadores = [
            "documento", "arquivo", "pdf", "imagem",
            "vejo", "analisei", "li", "observo",
            "conteúdo", "texto", "questão", "questões"
        ]
        
        # Verificar se menciona algo dos arquivos
        for anexo in anexos:
            nome_sem_ext = Path(anexo.nome).stem.lower()
            if nome_sem_ext in resposta_lower:
                return True
        
        # Verificar indicadores genéricos
        contagem = sum(1 for ind in indicadores if ind in resposta_lower)
        return contagem >= 2


# ============================================================
# INSTÂNCIAS GLOBAIS
# ============================================================

preparador_arquivos = PreparadorArquivos()
